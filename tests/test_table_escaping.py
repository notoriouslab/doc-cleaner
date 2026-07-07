"""Cross-format table-cell escaping tests (unify-table-cell-escaping).

Verifies the DOCX, PPTX, and XLSX paths apply the shared escape_cell rules:
pipes escaped, trailing backslash cannot eat a delimiter, multi-line cells
render on one physical line. (PDF is covered by tests/test_pdf_tables.py,
Numbers by tests/test_numbers_parser.py.)
"""
import os
import tempfile

import pytest


def _table_lines(md):
    return [ln for ln in md.split("\n") if ln.strip().startswith("|")]


def _col_count(line):
    return line.replace("\\\\", "").replace("\\|", "").count("|")


class TestDocxEscaping:
    @pytest.fixture
    def docx_path(self):
        docx_mod = pytest.importorskip("docx", reason="python-docx not installed")
        doc = docx_mod.Document()
        table = doc.add_table(rows=3, cols=2)
        table.rows[0].cells[0].text = "欄位"
        table.rows[0].cells[1].text = "值"
        table.rows[1].cells[0].text = "a|b"
        table.rows[1].cells[1].text = "C:\\"
        table.rows[2].cells[0].text = "第一行\n第二行"
        table.rows[2].cells[1].text = "x"
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc.save(f.name)
            yield f.name
        os.unlink(f.name)

    def test_special_cells_escaped(self, docx_path):
        from parsers.docx import parse
        out = parse(docx_path)
        assert "a\\|b" in out
        assert "C:\\\\" in out
        assert "第一行 第二行" in out
        lines = _table_lines(out)
        assert len({_col_count(ln) for ln in lines}) == 1  # constant columns


class TestPptxEscaping:
    @pytest.fixture
    def pptx_path(self):
        pptx_mod = pytest.importorskip("pptx", reason="python-pptx not installed")
        from pptx.util import Inches
        prs = pptx_mod.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        shape = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(6), Inches(3))
        tbl = shape.table
        tbl.cell(0, 0).text = "欄位"
        tbl.cell(0, 1).text = "值"
        tbl.cell(1, 0).text = "a|b"
        tbl.cell(1, 1).text = "C:\\"
        tbl.cell(2, 0).text = "第一行\n第二行"
        tbl.cell(2, 1).text = "x"
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            yield f.name
        os.unlink(f.name)

    def test_special_cells_escaped(self, pptx_path):
        from parsers.pptx import parse
        out = parse(pptx_path)
        assert "a\\|b" in out
        assert "C:\\\\" in out
        assert "第一行 第二行" in out
        lines = _table_lines(out)
        assert len({_col_count(ln) for ln in lines}) == 1
        # multi-line cell must not split the logical row
        assert len(lines) == 4  # header + separator + 2 data rows


class TestXlsxEscaping:
    @pytest.fixture
    def xlsx_path(self):
        openpyxl = pytest.importorskip("openpyxl", reason="openpyxl not installed")
        pytest.importorskip("pandas", reason="pandas not installed")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Price|USD", "說明", "金額"])          # header with pipe
        ws.append(["a|b", "第一行\n第二行", 100])
        ws.append(["C:\\", "純文字", 250])
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb.save(f.name)
            yield f.name
        os.unlink(f.name)

    @pytest.fixture
    def xlsx_control_path(self):
        openpyxl = pytest.importorskip("openpyxl", reason="openpyxl not installed")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["PriceUSD", "說明", "金額"])
        ws.append(["ab", "第一行", 100])
        ws.append(["C:", "純文字", 250])
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb.save(f.name)
            yield f.name
        os.unlink(f.name)

    def test_values_and_header_escaped_single_line(self, xlsx_path):
        from parsers.xlsx import parse
        out = parse(xlsx_path)
        assert "a\\|b" in out                    # data cell escaped
        assert "Price\\|USD" in out              # header cell escaped
        assert "C:\\\\" in out                   # trailing backslash escaped
        assert "第一行 第二行" in out            # newline collapsed → one physical row
        lines = _table_lines(out)
        assert len({_col_count(ln) for ln in lines}) == 1   # constant columns
        assert len(lines) == 4                   # header + separator + 2 data rows

    def test_numeric_column_unchanged(self, xlsx_path, xlsx_control_path):
        from parsers.xlsx import parse
        out = parse(xlsx_path)
        control = parse(xlsx_control_path)
        def nums(md):
            return [c.strip() for ln in _table_lines(md) for c in ln.split("|") if c.strip().isdigit()]
        assert nums(out) == nums(control) == ["100", "250"]
